import streamlit as st
from pptx import Presentation
import io
from PyPDF2 import PdfReader

def extract_text_from_pptx(file):
    prs = Presentation(file)
    slides_text = []
    for slide in prs.slides:
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):  # Verifica se o shape tem texto
                slide_text += shape.text + "\n"
        slides_text.append(slide_text.strip())
    return slides_text

def extract_text_from_pdf(file):
    reader = PdfReader(file)
    pages_text = [page.extract_text() for page in reader.pages if page.extract_text()]
    return pages_text

st.title("Upload e Extração de Texto de PPTX e PDF")

uploaded_file = st.file_uploader("Envie um arquivo .pptx ou .pdf", type=["pptx", "pdf"])

if uploaded_file is not None:
    file_type = uploaded_file.name.split(".")[-1]
    
    if file_type == "pptx":
        texts = extract_text_from_pptx(io.BytesIO(uploaded_file.read()))
    elif file_type == "pdf":
        texts = extract_text_from_pdf(io.BytesIO(uploaded_file.read()))
    
    for i, text in enumerate(texts):
        st.subheader(f"Página {i+1}")
        st.text_area(f"Texto da Página {i+1}", text, height=200)
    
    # Salvar em variável global (para uso posterior)
    st.session_state["extracted_texts"] = texts

    st.success("Texto extraído com sucesso!")
